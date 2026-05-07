"""
Génère le PowerPoint et le Word pour la présentation orale.
Lancer : python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from docx import Document
from docx.shared import Pt as DPt, RGBColor as DRGBColor, Inches as DInches
from docx.enum.text import WD_ALIGN_PARAGRAPH
import os

OUT_DIR = os.path.dirname(os.path.abspath(__file__))

# ─────────────────────────────────────────────────────────
# COULEURS
# ─────────────────────────────────────────────────────────
C_PRIMARY   = RGBColor(0xFF, 0x44, 0x58)   # rouge/rose principal
C_DARK      = RGBColor(0x1A, 0x1A, 0x2E)   # quasi-noir
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xF8, 0xF8, 0xFA)   # fond gris très clair
C_ACCENT    = RGBColor(0x22, 0xC5, 0x5E)   # vert match
C_GRAY      = RGBColor(0x6B, 0x72, 0x80)

# ─────────────────────────────────────────────────────────
# HELPERS PPTX
# ─────────────────────────────────────────────────────────
def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT,
                wrap=True):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def set_bg(slide, color: RGBColor):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bullets(slide, items, left, top, width, height,
                font_size=16, color=C_DARK, bullet_color=C_PRIMARY):
    txb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf  = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color


# ─────────────────────────────────────────────────────────
# DONNÉES DES SLIDES
# ─────────────────────────────────────────────────────────
SLIDES = [
    # (titre_slide, sous_titre, [bullets], note_orateur, dark_header)
    {
        "title": "Application de Rencontre",
        "subtitle": "Lohan Vignals — Master 1 POO\nProf. Adrien Escourrou — 2025/2026",
        "bullets": [],
        "note": "Bonjour. Je vais vous présenter mon projet de fin de semestre : une application mobile de rencontre, développée dans le cadre du cours de Programmation Orientée Objet. Le projet couvre l'ensemble du cycle de développement : conception, backend, frontend, et base de données.",
        "type": "title",
    },
    {
        "title": "Sommaire",
        "subtitle": "",
        "bullets": [
            "1.  Présentation du projet",
            "2.  Architecture technique",
            "3.  Fonctionnalités principales",
            "4.  Démonstration live",
            "5.  Concepts POO",
            "6.  Usage de l'IA",
            "7.  Difficultés rencontrées",
            "8.  Conclusion & perspectives",
        ],
        "note": "Je vais commencer par vous présenter le projet et son architecture, puis je détaillerai les fonctionnalités principales. On passera ensuite à une démo live de l'application. Je conclurai par les aspects POO, l'usage de l'IA, et les difficultés que j'ai rencontrées.",
        "type": "content",
    },
    {
        "title": "Présentation du projet",
        "subtitle": "Une appli de rencontre mobile complète",
        "bullets": [
            "✦  Swipe pour liker ou passer des profils",
            "✦  Matching mutuel → messagerie instantanée",
            "✦  Géolocalisation — trouver des profils proches",
            "✦  Profils enrichis : style de vie, type de relation, taille…",
            "✦  Sons, vibrations, notifications push",
            "",
            "Objectif : expérience fluide et réaliste d'une vraie application de rencontre.",
        ],
        "note": "L'idée est simple : reproduire l'expérience d'une vraie application de rencontre comme Tinder. L'utilisateur crée son profil, parcourt des profils proposés, peut les liker ou les passer. Quand deux personnes se likent mutuellement, c'est un match — et elles peuvent alors se parler.\n\nMais au-delà du swipe basique, j'ai voulu aller plus loin : géolocalisation pour trouver des profils proches, messagerie instantanée avec indicateurs de lecture, messages vocaux, filtres de recherche, et des profils enrichis avec le style de vie, le type de relation recherché, et la taille.",
        "type": "content",
    },
    {
        "title": "Architecture technique",
        "subtitle": "Stack fullstack — Frontend · Backend · Base de données",
        "bullets": [
            "FRONTEND          React Native / Expo SDK 54   →   iOS & Android",
            "BACKEND           FastAPI (Python 3.12)        →   API REST + WebSocket",
            "BASE DE DONNÉES   Supabase (PostgreSQL)        →   DB + Storage fichiers",
            "",
            "Communication :",
            "   • HTTP REST  — requêtes standard (profils, like, upload…)",
            "   • WebSocket  — temps réel (messages, statut en ligne, appels)",
        ],
        "note": "L'architecture est fullstack. Le frontend est une application React Native développée avec Expo, qui tourne aussi bien sur iOS qu'Android. Le backend est une API REST en Python avec FastAPI. La base de données est Supabase, qui est un service PostgreSQL hébergé, que j'utilise aussi pour le stockage des fichiers.\n\nLa communication se fait de deux façons : HTTP classique pour les requêtes standard, et WebSocket pour tout ce qui est temps réel — les messages, les indicateurs de frappe, les statuts en ligne.",
        "type": "content",
    },
    {
        "title": "Backend — FastAPI",
        "subtitle": "20+ endpoints REST · Services séparés · Modèles Pydantic",
        "bullets": [
            "Services organisés par responsabilité :",
            "   UserService         →  inscription, connexion, mise à jour profil",
            "   InteractionService  →  swipe, matching, cooldown dislike",
            "   PhotoService        →  upload vers Supabase Storage",
            "   MatchService        →  détection de match mutuel",
            "",
            "Modèles Pydantic → validation automatique des données",
            "Authentification → hachage bcrypt des mots de passe",
        ],
        "note": "Le backend expose une vingtaine d'endpoints REST. J'ai organisé le code en services séparés selon le principe de responsabilité unique : un service pour les utilisateurs, un pour les interactions comme le swipe et le matching, un pour les photos, un pour la logique de match.\n\nLes données entrantes sont validées automatiquement grâce aux modèles Pydantic. La gestion des mots de passe utilise bcrypt pour le hachage.",
        "type": "content",
    },
    {
        "title": "Frontend — React Native / Expo",
        "subtitle": "Application mobile iOS & Android",
        "bullets": [
            "Écrans principaux :",
            "   HomeScreen      →  swipe de profils",
            "   MatchesScreen   →  liste des matchs",
            "   ChatScreen      →  messagerie",
            "   ProfileScreen   →  édition du profil",
            "   MapScreen       →  carte géolocalisée",
            "   FiltersScreen   →  filtres de recherche",
            "",
            "Gestes natifs PanResponder · Animated API · expo-av / expo-haptics",
        ],
        "note": "Côté frontend, l'application est organisée en onglets. Le swipe est géré nativement avec PanResponder — l'API gesture de React Native — ce qui donne une sensation très fluide. Les animations sont gérées avec l'Animated API. Et pour renforcer le côté immersif, j'ai ajouté des sons et des vibrations à chaque interaction.",
        "type": "content",
    },
    {
        "title": "Base de données — Supabase",
        "subtitle": "PostgreSQL hébergé + Storage fichiers",
        "bullets": [
            "users      →  profil, coordonnées GPS, push_token",
            "photos     →  photos multiples par utilisateur",
            "likes      →  historique des swipes droits",
            "dislikes   →  cooldown 24h avant re-proposition",
            "matches    →  paires formées après match mutuel",
            "messages   →  texte, audio, réactions",
            "blocks     →  signalements / blocages",
            "",
            "Storage : photos de profil · messages vocaux",
        ],
        "note": "La base de données contient sept tables principales. La table users stocke les informations de profil, les coordonnées GPS et le token de notifications push. J'ai une table dislikes séparée pour gérer le cooldown de 24 heures — les profils repassés redeviennent visibles après 24 heures au lieu d'être cachés définitivement.",
        "type": "content",
    },
    {
        "title": "Fonctionnalité : Swipe & Matching",
        "subtitle": "Le cœur de l'application",
        "bullets": [
            "Geste natif gauche/droite  ou  boutons ❌ / ♥",
            "",
            "Algorithme de matching :",
            "   → A like B  +  B a déjà liké A  →  MATCH !",
            "   → Notification push envoyée immédiatement",
            "",
            "Cooldown dislike : 24 heures",
            "   → Les profils passés sont re-proposés après 24h",
            "",
            "Filtres : âge min/max · distance maximale (km)",
            "Distance calculée par formule Haversine (GPS)",
            "Coordonnées floutées ±800m pour la vie privée",
        ],
        "note": "Le cœur de l'application, c'est le swipe. Quand l'utilisateur swipe à droite, une requête est envoyée au backend qui vérifie si l'autre personne a déjà liké aussi. Si oui, c'est un match — le backend crée l'entrée en base et envoie une notification push.\n\nPour la distance, j'utilise la formule Haversine. Et pour préserver la vie privée, les coordonnées affichées sur la carte sont légèrement floues — décalées aléatoirement d'environ 800 mètres.",
        "type": "content",
    },
    {
        "title": "Fonctionnalité : Messagerie temps réel",
        "subtitle": "WebSocket + notifications push",
        "bullets": [
            "WebSocket persistent par utilisateur connecté",
            "",
            "Indicateurs temps réel :",
            "   ✓  « En train d'écrire... »",
            "   ✓  « Vu » — messages lus",
            "   ✓  Statut en ligne / hors ligne",
            "",
            "Messages vocaux (enregistrement + upload Supabase)",
            "Réactions aux messages (emojis)",
            "Appels vidéo via WebRTC (signaling WebSocket)",
            "Notifications push si l'utilisateur est hors ligne",
        ],
        "note": "La messagerie est entièrement temps réel via WebSocket. Chaque utilisateur connecté maintient une connexion persistante avec le serveur. Ça permet d'envoyer les messages instantanément, d'afficher l'indicateur de frappe, de montrer quand les messages sont lus, et de savoir si l'autre est en ligne. Si l'utilisateur est hors ligne, une notification push est envoyée via Expo.",
        "type": "content",
    },
    {
        "title": "Démonstration live",
        "subtitle": "",
        "bullets": [
            "Scénario :",
            "",
            "1.  HomeScreen → parcourir des profils",
            "2.  Swipe droit → son + vibration",
            "3.  Match → écran de match",
            "4.  Chat → envoyer un message",
        ],
        "note": "Je vais maintenant vous faire une démonstration live de l'application.\n\n[Montrer sur le téléphone]\n1. L'écran de swipe avec les profils, la carte et les filtres\n2. Un swipe à droite avec le son et la vibration\n3. L'écran de match\n4. Envoyer un message dans le chat",
        "type": "demo",
    },
    {
        "title": "Concepts POO",
        "subtitle": "Classes · Encapsulation · Polymorphisme",
        "bullets": [
            "Classes Pydantic — héritage de BaseModel :",
            "   class UserCreate(BaseModel):",
            "       username: str",
            "       date_of_birth: str",
            "       gender: Optional[str] = None",
            "",
            "   class LikeAction(BaseModel):",
            "       user_id: str",
            "       liked_user_id: str",
            "",
            "Encapsulation : chaque Service encapsule sa logique métier",
            "Polymorphisme : FastAPI route handlers polymorphes selon le type HTTP",
        ],
        "note": "Sur l'aspect Programmation Orientée Objet, le backend utilise des classes Pydantic pour modéliser les données. UserCreate hérite de BaseModel et définit les champs avec leur type et leur validation automatique. C'est le principe d'encapsulation : la logique de validation est dans la classe.\n\nL'architecture en services séparés applique le principe d'encapsulation à l'échelle des modules. Chaque service est responsable d'un seul domaine métier.",
        "type": "content",
    },
    {
        "title": "Usage de l'Intelligence Artificielle",
        "subtitle": "Claude (Anthropic) — outil d'accélération",
        "bullets": [
            "IA utilisée : Claude (Anthropic)",
            "",
            "Cas d'usage :",
            "   • Structure initiale des écrans React Native",
            "   • Algorithme Haversine (distance GPS)",
            "   • Logique WebSocket (messagerie temps réel)",
            "   • Composant SwipeCard (animations PanResponder)",
            "   • Gestion des sons et vibrations",
            "",
            "Tout le code IA est déclaré :",
            "   # ############### CODE IA (Claude) ###############",
            "",
            "Fait manuellement : architecture, schéma BDD, intégration Supabase, tests, UX",
        ],
        "note": "J'ai utilisé Claude d'Anthropic comme assistant de développement. L'IA m'a aidé à générer les composants complexes — notamment le SwipeCard avec ses animations, la logique WebSocket, et l'algorithme Haversine. Tout le code généré est explicitement déclaré avec les marqueurs requis dans le sujet.\n\nCe que j'ai fait manuellement : la conception de l'architecture, le schéma de la base de données, l'intégration avec Supabase, les tests sur appareil réel, et tous les ajustements d'UX après les tests. L'IA est un outil d'accélération, pas un remplaçant du raisonnement.",
        "type": "content",
    },
    {
        "title": "Difficultés rencontrées",
        "subtitle": "Problèmes → Solutions",
        "bullets": [
            "WebSocket mobile",
            "   Connexion coupée en arrière-plan",
            "   → Reconnexion automatique + ping keepalive",
            "",
            "DatePicker cross-platform",
            "   iOS et Android se comportent différemment",
            "   → Mode spinner (iOS) / dialog (Android) avec gestion séparée",
            "",
            "Performance du swipe",
            "   Animations saccadées avec beaucoup de profils",
            "   → useNativeDriver + préchargement de la carte suivante",
        ],
        "note": "Trois difficultés principales. D'abord, le WebSocket sur mobile : les connexions se coupent quand l'app passe en arrière-plan, j'ai dû gérer la reconnexion automatique. Ensuite, le date picker qui fonctionne différemment sur iOS et Android. Et enfin la performance du swipe : pour que les animations soient fluides, il faut utiliser le native driver et précharger la carte suivante.",
        "type": "content",
    },
    {
        "title": "Conclusion & Perspectives",
        "subtitle": "",
        "bullets": [
            "Ce projet m'a permis de :",
            "   ✓  Concevoir une architecture fullstack complète",
            "   ✓  Travailler avec des APIs temps réel (WebSocket)",
            "   ✓  Gérer la géolocalisation et les notifications push",
            "   ✓  Utiliser l'IA comme outil de développement",
            "",
            "Améliorations possibles :",
            "   →  Algorithme de recommandation (score de compatibilité)",
            "   →  Vérification d'identité (photo selfie)",
            "   →  Abonnement premium (boosts, super likes)",
            "   →  Modération automatique des photos (IA)",
            "",
            "Merci — Questions ?",
        ],
        "note": "Pour conclure, ce projet m'a permis de développer une application complète de bout en bout, avec des fonctionnalités réellement complexes — WebSocket, géolocalisation, notifications push, appels vidéo. Les axes d'amélioration sont nombreux : un algorithme de recommandation, une vérification d'identité, ou un système d'abonnement premium.\n\nMerci pour votre attention. Je suis disponible pour vos questions.",
        "type": "conclusion",
    },
]

# ─────────────────────────────────────────────────────────
# GÉNÉRATION POWERPOINT
# ─────────────────────────────────────────────────────────
def build_pptx():
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for i, s in enumerate(SLIDES):
        slide = prs.slides.add_slide(blank_layout)

        if s["type"] == "title":
            set_bg(slide, C_DARK)
            # Bande colorée gauche
            add_rect(slide, 0, 0, 0.6, 7.5, C_PRIMARY)
            # Numéro de slide
            add_textbox(slide, f"{i+1}/{len(SLIDES)}", 11.5, 6.8, 1.5, 0.5,
                        font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)
            # Titre
            add_textbox(slide, s["title"], 1.0, 2.2, 11.0, 1.4,
                        font_size=44, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
            # Ligne décorative
            add_rect(slide, 1.0, 3.75, 3.0, 0.06, C_PRIMARY)
            # Sous-titre
            add_textbox(slide, s["subtitle"], 1.0, 3.95, 11.0, 1.2,
                        font_size=18, color=C_LIGHT, align=PP_ALIGN.LEFT)

        elif s["type"] == "demo":
            set_bg(slide, C_PRIMARY)
            add_textbox(slide, s["title"], 0.5, 2.8, 12.3, 1.5,
                        font_size=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
            add_textbox(slide, f"{i+1}/{len(SLIDES)}", 11.5, 6.8, 1.5, 0.5,
                        font_size=10, color=C_WHITE, align=PP_ALIGN.RIGHT)
            if s["bullets"]:
                add_bullets(slide, s["bullets"], 3.5, 4.5, 6.0, 2.5,
                            font_size=15, color=C_WHITE)

        elif s["type"] == "conclusion":
            set_bg(slide, C_DARK)
            add_rect(slide, 0, 0, 0.6, 7.5, C_PRIMARY)
            add_textbox(slide, s["title"], 1.0, 0.4, 11.0, 0.9,
                        font_size=28, bold=True, color=C_WHITE)
            add_rect(slide, 1.0, 1.35, 2.5, 0.05, C_PRIMARY)
            add_textbox(slide, f"{i+1}/{len(SLIDES)}", 11.5, 6.8, 1.5, 0.5,
                        font_size=10, color=C_GRAY, align=PP_ALIGN.RIGHT)
            add_bullets(slide, s["bullets"], 1.0, 1.5, 11.0, 5.5,
                        font_size=15, color=C_WHITE)

        else:
            set_bg(slide, C_LIGHT)
            # Header bar
            add_rect(slide, 0, 0, 13.33, 1.35, C_DARK)
            # Bande rouge gauche header
            add_rect(slide, 0, 0, 0.35, 1.35, C_PRIMARY)
            # Titre en blanc
            add_textbox(slide, s["title"], 0.5, 0.12, 10.0, 0.72,
                        font_size=24, bold=True, color=C_WHITE)
            # Numéro
            add_textbox(slide, f"{i+1}/{len(SLIDES)}", 11.5, 0.18, 1.6, 0.6,
                        font_size=12, color=C_GRAY, align=PP_ALIGN.RIGHT)
            # Sous-titre
            if s["subtitle"]:
                add_textbox(slide, s["subtitle"], 0.5, 0.78, 12.0, 0.45,
                            font_size=13, color=C_PRIMARY)
            # Bullets
            if s["bullets"]:
                add_bullets(slide, s["bullets"], 0.7, 1.5, 11.9, 5.7,
                            font_size=15, color=C_DARK)

        # Notes orateur
        if s["note"]:
            slide.notes_slide.notes_text_frame.text = s["note"]

    path = os.path.join(OUT_DIR, "Presentation_Application_Rencontre.pptx")
    prs.save(path)
    print(f"OK PowerPoint cree : {path}")


# ─────────────────────────────────────────────────────────
# GÉNÉRATION WORD (SCRIPT)
# ─────────────────────────────────────────────────────────
def build_docx():
    doc = Document()

    # Marges
    for section in doc.sections:
        section.top_margin    = DInches(1)
        section.bottom_margin = DInches(1)
        section.left_margin   = DInches(1.2)
        section.right_margin  = DInches(1.2)

    # Titre principal
    t = doc.add_heading("Script de présentation orale", level=0)
    t.alignment = WD_ALIGN_PARAGRAPH.CENTER

    sub = doc.add_paragraph("Application de Rencontre — Lohan Vignals — Master 1 POO")
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sub.runs[0].font.size = DPt(13)
    sub.runs[0].font.color.rgb = DRGBColor(0x6B, 0x72, 0x80)

    doc.add_paragraph()

    TIMING = [
        "30 sec", "20 sec", "1 min 30", "2 min", "1 min 30",
        "1 min", "1 min", "1 min 30", "1 min", "2 min (démo)",
        "1 min 30", "1 min 30", "1 min", "1 min",
    ]

    for i, s in enumerate(SLIDES):
        # ── En-tête slide ──
        hdr = doc.add_heading(f"Slide {i+1} — {s['title']}", level=2)
        hdr.runs[0].font.color.rgb = DRGBColor(0xFF, 0x44, 0x58)

        timing = doc.add_paragraph(f"⏱  Durée estimée : {TIMING[i]}")
        timing.runs[0].font.size  = DPt(10)
        timing.runs[0].font.color.rgb = DRGBColor(0x6B, 0x72, 0x80)
        timing.runs[0].font.italic = True

        # ── Script ──
        if s["note"]:
            for line in s["note"].split("\n"):
                if line.startswith("[") or line.startswith("("):
                    p = doc.add_paragraph(line)
                    p.runs[0].font.italic = True
                    p.runs[0].font.color.rgb = DRGBColor(0x6B, 0x72, 0x80)
                elif line.strip() == "":
                    doc.add_paragraph()
                else:
                    p = doc.add_paragraph(f"« {line} »")
                    p.runs[0].font.size = DPt(12)

        doc.add_paragraph()
        doc.add_paragraph("─" * 60)
        doc.add_paragraph()

    # ── Questions types ──
    doc.add_heading("Questions types du professeur", level=1)

    QA = [
        (
            "Pourquoi FastAPI plutôt que Flask ou Django ?",
            "FastAPI est asynchrone nativement, ce qui est essentiel pour les WebSockets. Il génère aussi automatiquement la documentation de l'API. Flask aurait nécessité plus de bibliothèques externes."
        ),
        (
            "Comment tu gères la sécurité ?",
            "Les mots de passe sont hachés avec bcrypt avant stockage. Les coordonnées GPS sont floutées pour la vie privée. Le système de blocage empêche les contacts non désirés."
        ),
        (
            "C'est quoi le polymorphisme dans ton code ?",
            "Les classes Pydantic (UserCreate, LikeAction, UserUpdate...) héritent toutes de BaseModel et redéfinissent leurs champs — c'est du polymorphisme de sous-type. FastAPI utilise ces classes de façon polymorphe pour valider n'importe quel type de requête."
        ),
        (
            "Tu aurais fait quoi différemment ?",
            "J'aurais intégré une vraie authentification JWT dès le départ. Et j'aurais écrit des tests unitaires en parallèle du développement plutôt qu'à la fin."
        ),
        (
            "Pourquoi Supabase et pas une autre base de données ?",
            "Supabase combine PostgreSQL hébergé, le stockage de fichiers, et une API REST générée automatiquement. C'est un service tout-en-un qui m'a évité de configurer une infrastructure complète."
        ),
        (
            "Comment fonctionne le matching ?",
            "Quand un utilisateur like quelqu'un, le backend vérifie si l'autre a déjà liké en retour. Si oui, un match est créé et les deux utilisateurs reçoivent une notification push."
        ),
    ]

    for q, a in QA:
        doc.add_heading(f"Q : {q}", level=3)
        p = doc.add_paragraph(f"R : {a}")
        p.runs[0].font.size = DPt(12)
        doc.add_paragraph()

    path = os.path.join(OUT_DIR, "Script_Presentation_Rencontre.docx")
    doc.save(path)
    print(f"OK Word cree : {path}")


if __name__ == "__main__":
    build_pptx()
    build_docx()
    print("\nFichiers generes dans :", OUT_DIR)
